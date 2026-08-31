import os
import re
from typing import List, Dict, Any, Tuple
import aiofiles
from sqlalchemy.orm import Session
from pypdf import PdfReader
import docx
from pptx import Presentation
from app.models.document import Document, DocumentChunk
from app.ai.providers.embedding_provider import get_embedding_provider
from app.core.logging import logger


class DocumentIngestionService:
    def __init__(self, db: Session):
        self.db = db
        self.embedding_provider = get_embedding_provider()

    async def process_and_index_document(self, document_id: str) -> bool:
        doc = self.db.query(Document).filter(Document.id == document_id).first()
        if not doc:
            logger.error(f"Document {document_id} not found for indexing.")
            return False

        try:
            doc.processing_status = "processing"
            self.db.commit()

            # 1. Extract text and structural pages
            pages_data = self._extract_pages(doc.storage_path, doc.file_type)
            doc.page_count = len(pages_data)

            # 2. Hierarchical Chunking
            chunks_to_create = []
            for page_num, section_title, page_text in pages_data:
                page_chunks = self._chunk_text(page_text, max_chunk_size=400, overlap=50)
                for chunk_text in page_chunks:
                    chunks_to_create.append((page_num, section_title, chunk_text))

            # 3. Vector Embeddings
            raw_texts = [c[2] for c in chunks_to_create]
            embeddings = await self.embedding_provider.embed_documents(raw_texts)

            # 4. Save chunks to database
            for (page_num, section_title, chunk_text), emb in zip(chunks_to_create, embeddings):
                db_chunk = DocumentChunk(
                    document_id=doc.id,
                    chunk_text=chunk_text,
                    page_number=page_num,
                    section_title=section_title,
                    embedding=emb,
                    chunk_metadata={"filename": doc.filename, "char_len": len(chunk_text)}
                )
                self.db.add(db_chunk)

            doc.processing_status = "indexed"
            self.db.commit()
            logger.info(f"Successfully indexed document {doc.filename} with {len(chunks_to_create)} chunks.")
            return True

        except Exception as e:
            logger.error(f"Error processing document {document_id}: {e}", exc_info=True)
            doc.processing_status = "failed"
            doc.doc_metadata = {"error": str(e)}
            self.db.commit()
            return False

    def _extract_pages(self, file_path: str, file_type: str) -> List[Tuple[int, str, str]]:
        """
        Extracts structured (page_number, section_title, text) tuples.
        """
        pages = []
        ext = file_type.lower().lstrip(".")

        if ext == "pdf":
            reader = PdfReader(file_path)
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                # Attempt to extract heading from first non-empty line
                lines = [l.strip() for l in text.split("\n") if l.strip()]
                section = lines[0][:60] if lines else f"Page {i+1}"
                if text.strip():
                    pages.append((i + 1, section, text))

        elif ext in ["docx", "doc"]:
            doc_obj = docx.Document(file_path)
            full_text = []
            current_section = "Main Content"
            for para in doc_obj.paragraphs:
                if para.style.name.startswith("Heading"):
                    current_section = para.text.strip() or current_section
                if para.text.strip():
                    full_text.append(para.text.strip())
            joined = "\n\n".join(full_text)
            pages.append((1, current_section, joined))

        elif ext in ["pptx", "ppt"]:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                slide_title = f"Slide {i+1}"
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_texts.append(paragraph.text.strip())
                if slide_texts:
                    slide_title = slide_texts[0][:50]
                joined = "\n".join(slide_texts)
                pages.append((i + 1, slide_title, joined))

        else:  # txt / md / fallback
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                # Split roughly by chapters or markdown headers
                sections = re.split(r'\n(?=#+ )', content)
                for i, sec in enumerate(sections):
                    lines = sec.strip().split("\n")
                    sec_title = lines[0].replace("#", "").strip()[:50] if lines else f"Section {i+1}"
                    pages.append((i + 1, sec_title, sec))

        return pages

    def _chunk_text(self, text: str, max_chunk_size: int = 400, overlap: int = 50) -> List[str]:
        words = text.split()
        if not words:
            return []
        
        chunks = []
        start = 0
        while start < len(words):
            end = min(start + max_chunk_size, len(words))
            chunk = " ".join(words[start:end])
            if chunk.strip():
                chunks.append(chunk)
            if end >= len(words):
                break
            start += max_chunk_size - overlap
        return chunks
